import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from models.ppt_state import (
    Image,
    PPTState,
    Position,
    Size,
    Slide,
    SlideSize,
    TextBox,
    TextStyle,
)
from utils.emu import emu_to_px

MAX_FILE_SIZE = 50 * 1024 * 1024  # 50 MB
ALLOWED_MIME_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/octet-stream",
}

IMAGE_PLACEHOLDER_TYPES = {
    PP_PLACEHOLDER.PICTURE,
    PP_PLACEHOLDER.MEDIA_CLIP,
    PP_PLACEHOLDER.OBJECT,
}


def _validate_pptx(file_path: Path) -> None:
    """Validate file is a valid PPTX."""
    if not file_path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")
    if file_path.stat().st_size > MAX_FILE_SIZE:
        raise ValueError(f"File exceeds {MAX_FILE_SIZE} bytes limit")
    if not zipfile.is_zipfile(file_path):
        raise ValueError("File is not a valid ZIP/PPTX format")
    with zipfile.ZipFile(file_path, "r") as zf:
        if "ppt/presentation.xml" not in zf.namelist():
            raise ValueError("Missing required presentation.xml entry")


def _extract_textboxes(slide) -> list[TextBox]:
    """Extract text boxes from a slide."""
    textboxes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.is_placeholder and not shape.text_frame.text.strip():
            continue

        paragraphs = []
        for para in shape.text_frame.paragraphs:
            para_text = "".join(run.text for run in para.runs)
            paragraphs.append(para_text)
        content = "\n".join(paragraphs)

        # Extract style from first run
        style = TextStyle()
        if shape.text_frame.paragraphs:
            first_para = shape.text_frame.paragraphs[0]
            if first_para.runs:
                first_run = first_para.runs[0]
                font = first_run.font
                if font.size:
                    style.font_size_pt = font.size.pt
                try:
                    if font.color and font.color.rgb:
                        style.font_color = f"#{font.color.rgb}"
                except AttributeError:
                    pass
                style.bold = font.bold
                style.italic = font.italic

        textboxes.append(
            TextBox(
                content=content,
                position=Position(
                    x_emu=shape.left,
                    y_emu=shape.top,
                    x_px=emu_to_px(shape.left),
                    y_px=emu_to_px(shape.top),
                ),
                size=Size(
                    width_emu=shape.width,
                    height_emu=shape.height,
                    width_px=emu_to_px(shape.width),
                    height_px=emu_to_px(shape.height),
                ),
                style=style,
            )
        )
    return textboxes


def _extract_images(slide) -> list[Image]:
    """Extract image placeholders from a slide."""
    images = []
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        if shape.placeholder_format.type not in IMAGE_PLACEHOLDER_TYPES:
            continue
        images.append(
            Image(
                position=Position(
                    x_emu=shape.left,
                    y_emu=shape.top,
                    x_px=emu_to_px(shape.left),
                    y_px=emu_to_px(shape.top),
                ),
                size=Size(
                    width_emu=shape.width,
                    height_emu=shape.height,
                    width_px=emu_to_px(shape.width),
                    height_px=emu_to_px(shape.height),
                ),
                placeholder_type=shape.placeholder_format.type.name.lower(),
            )
        )
    return images


def parse_pptx(file_path: str | Path, page_nums: list[int] | None = None) -> PPTState:
    """Parse a PPTX file into PPTState.

    Args:
        file_path: Path to the .pptx file.
        page_nums: Optional list of 1-based page numbers to extract.
            Defaults to first 50 pages.

    Returns:
        PPTState representing the parsed slides.
    """
    file_path = Path(file_path)
    _validate_pptx(file_path)

    prs = Presentation(str(file_path))
    total_slides = len(prs.slides)

    if page_nums is None:
        page_nums = list(range(1, min(total_slides + 1, 51)))  # default first 50
    else:
        page_nums = sorted(set(page_nums))
        if any(p < 1 or p > total_slides for p in page_nums):
            raise ValueError(f"Page numbers out of range (1-{total_slides})")
        if len(page_nums) > 50:
            raise ValueError("Supports at most 50 pages")

    slides = []
    for page_num in page_nums:
        slide = prs.slides[page_num - 1]
        textboxes = _extract_textboxes(slide)
        images = _extract_images(slide)
        slides.append(
            Slide(
                page_num=page_num,
                size=SlideSize(
                    width_emu=prs.slide_width,
                    height_emu=prs.slide_height,
                    width_px=emu_to_px(prs.slide_width),
                    height_px=emu_to_px(prs.slide_height),
                ),
                elements=textboxes + images,
            )
        )

    return PPTState(
        source_file=str(file_path),
        slide_count=len(slides),
        global_props=SlideSize(
            width_emu=prs.slide_width,
            height_emu=prs.slide_height,
            width_px=emu_to_px(prs.slide_width),
            height_px=emu_to_px(prs.slide_height),
        ),
        slides=slides,
    )
