import shutil
import tempfile
from pathlib import Path

from pptx import Presentation

from models.ppt_state import Image, PPTState, TextBox


def _replace_text_preserving_format(shape, new_content: str) -> None:
    """Replace shape text while preserving original formatting."""
    if not shape.has_text_frame:
        return
    text_frame = shape.text_frame
    paragraphs = text_frame.paragraphs
    if not paragraphs:
        return

    first_para = paragraphs[0]
    if not first_para.runs:
        run = first_para.add_run()
        run.text = new_content
        return

    first_run = first_para.runs[0]
    # Clear other paragraphs
    for para in paragraphs[1:]:
        para.clear()
    # Clear other runs in first paragraph
    for run in first_para.runs[1:]:
        run.text = ""

    first_run.text = new_content


def _find_shape_by_geometry(slide, left: int, top: int, width: int, height: int):
    """Find a shape by its geometry (position + size)."""
    for shape in slide.shapes:
        if (shape.left == left and shape.top == top and
                shape.width == width and shape.height == height):
            return shape
    return None


def _write_text_changes(slide, elements: list[TextBox]) -> None:
    """Apply text box changes to a slide."""
    for elem in elements:
        if elem.element_type != "textbox":
            continue
        shape = _find_shape_by_geometry(
            slide,
            left=elem.position.x_emu,
            top=elem.position.y_emu,
            width=elem.size.width_emu,
            height=elem.size.height_emu,
        )
        if shape and shape.has_text_frame:
            _replace_text_preserving_format(shape, elem.content)


def _write_image_changes(slide, elements: list[Image]) -> None:
    """Apply image placeholder changes to a slide."""
    # MVP: images are placeholders only; no binary replacement yet
    pass


def recompose_pptx(
    original_path: str | Path,
    ppt_state: PPTState,
    output_path: str | Path,
) -> Path:
    """Recompose a PPTX from PPTState, preserving original formatting.

    Args:
        original_path: Path to the original .pptx template.
        ppt_state: Modified PPTState with changes to apply.
        output_path: Destination path for the output .pptx.

    Returns:
        Path to the output file.
    """
    original_path = Path(original_path)
    output_path = Path(output_path)

    with tempfile.TemporaryDirectory() as tmp_dir:
        working_copy = Path(tmp_dir) / original_path.name
        shutil.copy2(original_path, working_copy)

        prs = Presentation(str(working_copy))

        for slide_state in ppt_state.slides:
            if slide_state.page_num < 1 or slide_state.page_num > len(prs.slides):
                continue
            slide = prs.slides[slide_state.page_num - 1]
            text_elems = [e for e in slide_state.elements if e.element_type == "textbox"]
            image_elems = [e for e in slide_state.elements if e.element_type == "image"]
            _write_text_changes(slide, text_elems)
            _write_image_changes(slide, image_elems)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

    return output_path
