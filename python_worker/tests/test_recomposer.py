from pathlib import Path

import pytest
from models.ppt_state import PPTState
from pptx import Presentation
from services.parser import parse_pptx
from services.recomposer import recompose_pptx

FIXTURES_DIR = Path(__file__).parent / "fixtures"


def test_recompose_no_changes():
    """Recomposing without changes should produce a valid PPTX."""
    fixture_path = FIXTURES_DIR / "sample.pptx"
    if not fixture_path.exists():
        pytest.skip("sample.pptx fixture not found")

    state = parse_pptx(str(fixture_path))
    output_path = FIXTURES_DIR / "output_no_changes.pptx"

    try:
        recompose_pptx(str(fixture_path), state, str(output_path))
        assert output_path.exists()
        # Should be parseable again
        state2 = parse_pptx(str(output_path))
        assert state2.source_file == "output_no_changes.pptx"
        assert state2.slide_count == state.slide_count
    finally:
        output_path.unlink(missing_ok=True)


def test_recompose_text_change():
    """Recomposing with a text change should update the content."""
    fixture_path = FIXTURES_DIR / "sample.pptx"
    if not fixture_path.exists():
        pytest.skip("sample.pptx fixture not found")

    state = parse_pptx(str(fixture_path))
    output_path = FIXTURES_DIR / "output_text_change.pptx"

    try:
        # Modify the first text box content
        if state.slides and state.slides[0].elements:
            text_elems = [e for e in state.slides[0].elements if e.element_type == "textbox"]
            if text_elems:
                text_elems[0].content = "Modified by PPT Agent"

        recompose_pptx(str(fixture_path), state, str(output_path))
        assert output_path.exists()

        state2 = parse_pptx(str(output_path))
        text_elems2 = [e for e in state2.slides[0].elements if e.element_type == "textbox"]
        assert text_elems2[0].content == "Modified by PPT Agent"
    finally:
        output_path.unlink(missing_ok=True)


def test_recompose_text_color_change_is_written_to_pptx():
    """Recomposing with a text color change should update the PPTX run color."""
    fixture_path = FIXTURES_DIR / "sample.pptx"
    if not fixture_path.exists():
        pytest.skip("sample.pptx fixture not found")

    state = parse_pptx(str(fixture_path))
    output_path = FIXTURES_DIR / "output_text_color_change.pptx"

    try:
        text_elems = [
            e
            for e in state.slides[0].elements
            if e.element_type == "textbox"
        ]
        if not text_elems:
            pytest.skip("sample.pptx fixture has no text boxes")
        text_elems[0].style.font_color = "#0000FF"

        recompose_pptx(str(fixture_path), state, str(output_path))
        assert output_path.exists()

        prs = Presentation(str(output_path))
        shape = next(shape for shape in prs.slides[0].shapes if shape.has_text_frame)
        run = shape.text_frame.paragraphs[0].runs[0]
        assert str(run.font.color.rgb) == "0000FF"
    finally:
        output_path.unlink(missing_ok=True)
